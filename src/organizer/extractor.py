# =============================================================================
# Local File Organizer - Extractor Component
# =============================================================================
"""
Extractor: Content extraction from various file formats.

The Extractor is the second stage of the Local-First pipeline:
1. Detects MIME types
2. Extracts text content based on file type
3. Truncates content to maximum excerpt size
4. Enriches FileRecord with content_excerpt and mime fields

Supported Formats:
- Text: .txt, .md, .csv, .json, .xml, .html, .log
- PDF: .pdf (via pdfplumber, first N pages)
- Word: .docx (via python-docx)
- PowerPoint: .pptx (via python-pptx)
- Excel: .xlsx, .xls (via pandas)
- Images: .jpg, .png, etc. (metadata only via Pillow)

Safety Features:
- Never modifies source files
- Graceful handling of corrupted/encrypted files
- Tracks statistics for audit/debugging
"""
import mimetypes
from pathlib import Path
from typing import Generator, Optional, List
import logging

from src.organizer.models import FileRecord


# Initialize mimetypes database
mimetypes.init()

# Configure logging
logger = logging.getLogger(__name__)


# =============================================================================
# Constants
# =============================================================================

DEFAULT_MAX_EXCERPT_BYTES: int = 8192  # 8KB
DEFAULT_MAX_PDF_PAGES: int = 5

# Supported text extensions (can be read directly)
TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst",
    ".csv", ".tsv",
    ".json", ".yaml", ".yml",
    ".xml", ".html", ".htm",
    ".log", ".ini", ".cfg", ".conf",
    ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".hpp",
    ".css", ".scss", ".sass",
    ".sh", ".bash", ".zsh",
    ".sql",
}

# Image extensions (metadata extraction only)
IMAGE_EXTENSIONS = {
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".tif",
    ".webp", ".heic", ".heif", ".raw", ".cr2", ".nef", ".arw",
}

# Audio extensions
AUDIO_EXTENSIONS = {
    ".mp3", ".wav", ".flac", ".ogg", ".aac", ".m4a", ".wma",
    ".opus", ".aiff", ".alac",
}

# Video extensions
VIDEO_EXTENSIONS = {
    ".mp4", ".avi", ".mkv", ".mov", ".wmv", ".flv", ".webm",
    ".m4v", ".mpeg", ".mpg", ".3gp",
}


# =============================================================================
# Lazy Imports (optional dependencies)
# =============================================================================

def _get_pdfplumber():
    """Lazy import pdfplumber."""
    try:
        import pdfplumber
        return pdfplumber
    except ImportError:
        logger.warning("pdfplumber not installed. PDF extraction disabled.")
        return None


def _get_docx():
    """Lazy import python-docx."""
    try:
        from docx import Document
        return Document
    except ImportError:
        logger.warning("python-docx not installed. DOCX extraction disabled.")
        return None


def _get_pptx():
    """Lazy import python-pptx."""
    try:
        from pptx import Presentation
        return Presentation
    except ImportError:
        logger.warning("python-pptx not installed. PPTX extraction disabled.")
        return None


def _get_pandas():
    """Lazy import pandas."""
    try:
        import pandas as pd
        return pd
    except ImportError:
        logger.warning("pandas not installed. Excel extraction disabled.")
        return None


def _get_pillow():
    """Lazy import Pillow."""
    try:
        from PIL import Image
        return Image
    except ImportError:
        logger.warning("Pillow not installed. Image metadata extraction disabled.")
        return None


def _get_mutagen():
    """Lazy import mutagen for audio metadata."""
    try:
        import mutagen
        return mutagen
    except ImportError:
        logger.warning("mutagen not installed. Audio metadata extraction disabled.")
        return None


def _get_ffprobe():
    """Check if ffprobe is available for video metadata."""
    import shutil
    import subprocess
    
    ffprobe_path = shutil.which("ffprobe")
    if ffprobe_path:
        return ffprobe_path
    
    # Try common locations on Windows
    common_paths = [
        r"C:\ffmpeg\bin\ffprobe.exe",
        r"C:\Program Files\ffmpeg\bin\ffprobe.exe",
        r"C:\Program Files (x86)\ffmpeg\bin\ffprobe.exe",
    ]
    for path in common_paths:
        if Path(path).exists():
            return path
    
    logger.warning("ffprobe not found. Video metadata extraction disabled.")
    return None


# Module-level lazy imports for mocking in tests
pdfplumber = None
Document = None
Presentation = None
pd = None
Image = None
mutagen = None
ffprobe_path = None


def _init_lazy_imports():
    """Initialize lazy imports on first use."""
    global pdfplumber, Document, Presentation, pd, Image, mutagen, ffprobe_path
    if pdfplumber is None:
        pdfplumber = _get_pdfplumber()
    if Document is None:
        Document = _get_docx()
    if Presentation is None:
        Presentation = _get_pptx()
    if pd is None:
        pd = _get_pandas()
    if Image is None:
        Image = _get_pillow()
    if mutagen is None:
        mutagen = _get_mutagen()
    if ffprobe_path is None:
        ffprobe_path = _get_ffprobe()


# =============================================================================
# Helper Functions
# =============================================================================

def detect_mime_type(file_path: Path) -> Optional[str]:
    """
    Detect MIME type of a file.

    Uses mimetypes module with extension-based fallback.

    Args:
        file_path: Path to the file

    Returns:
        MIME type string or None if detection fails
    """
    # Try mimetypes first (extension-based)
    mime_type, _ = mimetypes.guess_type(str(file_path))
    
    if mime_type:
        return mime_type
    
    # Fallback to generic binary
    return "application/octet-stream"


def truncate_content(content: str, max_bytes: int = DEFAULT_MAX_EXCERPT_BYTES) -> str:
    """
    Truncate content to maximum byte size.

    Handles multi-byte Unicode characters safely.

    Args:
        content: Content to truncate
        max_bytes: Maximum size in bytes

    Returns:
        Truncated content with marker if truncated
    """
    if content is None:
        return None
    
    encoded = content.encode("utf-8")
    
    if len(encoded) <= max_bytes:
        return content
    
    # Truncate at byte boundary, avoiding splitting multi-byte chars
    truncated = encoded[:max_bytes]
    
    # Decode safely, replacing incomplete chars
    try:
        result = truncated.decode("utf-8", errors="ignore")
    except UnicodeDecodeError:
        result = truncated.decode("utf-8", errors="replace")
    
    # Add truncation marker
    return result.rstrip() + "\n[TRUNCATED]"


def extract_text_content(
    file_path: Path,
    max_bytes: int = DEFAULT_MAX_EXCERPT_BYTES
) -> Optional[str]:
    """
    Extract content from plain text files.

    Args:
        file_path: Path to the text file
        max_bytes: Maximum content size

    Returns:
        Extracted text content or None on error
    """
    if not file_path.exists():
        return None
    
    try:
        # Try UTF-8 first
        try:
            content = file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            # Fallback to latin-1 (accepts any byte)
            content = file_path.read_text(encoding="latin-1")
        
        return truncate_content(content, max_bytes)
    except Exception as e:
        logger.warning(f"Failed to read text file {file_path}: {e}")
        return None


def extract_pdf_content(
    file_path: Path,
    max_pages: int = DEFAULT_MAX_PDF_PAGES,
    max_bytes: int = DEFAULT_MAX_EXCERPT_BYTES
) -> Optional[str]:
    """
    Extract text content from PDF files.

    Uses pdfplumber for extraction, limiting to first N pages.

    Args:
        file_path: Path to the PDF file
        max_pages: Maximum number of pages to extract
        max_bytes: Maximum content size

    Returns:
        Extracted text content or None on error
    """
    _init_lazy_imports()
    
    if pdfplumber is None or not file_path.exists():
        return None
    
    try:
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages[:max_pages]):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"[Page {i + 1}]\n{page_text}")
        
        if not text_parts:
            return None
        
        content = "\n\n".join(text_parts)
        return truncate_content(content, max_bytes)
    except Exception as e:
        logger.warning(f"Failed to extract PDF content from {file_path}: {e}")
        return None


def extract_docx_content(
    file_path: Path,
    max_bytes: int = DEFAULT_MAX_EXCERPT_BYTES
) -> Optional[str]:
    """
    Extract text content from DOCX files.

    Uses python-docx for extraction.

    Args:
        file_path: Path to the DOCX file
        max_bytes: Maximum content size

    Returns:
        Extracted text content or None on error
    """
    _init_lazy_imports()
    
    if Document is None or not file_path.exists():
        return None
    
    try:
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        
        if not paragraphs:
            return None
        
        content = "\n\n".join(paragraphs)
        return truncate_content(content, max_bytes)
    except Exception as e:
        logger.warning(f"Failed to extract DOCX content from {file_path}: {e}")
        return None


def extract_pptx_content(
    file_path: Path,
    max_bytes: int = DEFAULT_MAX_EXCERPT_BYTES
) -> Optional[str]:
    """
    Extract text content from PPTX files.

    Uses python-pptx for extraction, getting text from all shapes.

    Args:
        file_path: Path to the PPTX file
        max_bytes: Maximum content size

    Returns:
        Extracted text content or None on error
    """
    _init_lazy_imports()
    
    if Presentation is None or not file_path.exists():
        return None
    
    try:
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text)
            
            if slide_texts:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
        
        if not text_parts:
            return None
        
        content = "\n\n".join(text_parts)
        return truncate_content(content, max_bytes)
    except Exception as e:
        logger.warning(f"Failed to extract PPTX content from {file_path}: {e}")
        return None


def extract_xlsx_content(
    file_path: Path,
    max_rows: int = 10,
    max_bytes: int = DEFAULT_MAX_EXCERPT_BYTES
) -> Optional[str]:
    """
    Extract content from Excel files.

    Uses pandas for extraction, getting sheet names and first N rows.

    Args:
        file_path: Path to the Excel file
        max_rows: Maximum rows per sheet to extract
        max_bytes: Maximum content size

    Returns:
        Extracted content or None on error
    """
    _init_lazy_imports()
    
    if pd is None or not file_path.exists():
        return None
    
    try:
        text_parts = []
        
        with pd.ExcelFile(file_path) as xl:
            for sheet_name in xl.sheet_names[:5]:  # Max 5 sheets
                df = pd.read_excel(xl, sheet_name=sheet_name, nrows=max_rows)
                
                if df.empty:
                    continue
                
                # Get column names
                columns = df.columns.tolist()
                # Get preview of data
                preview = df.head(max_rows).to_string(index=False)
                
                text_parts.append(
                    f"[Sheet: {sheet_name}]\n"
                    f"Columns: {', '.join(str(c) for c in columns)}\n"
                    f"{preview}"
                )
        
        if not text_parts:
            return None
        
        content = "\n\n".join(text_parts)
        return truncate_content(content, max_bytes)
    except Exception as e:
        logger.warning(f"Failed to extract Excel content from {file_path}: {e}")
        return None


def extract_image_metadata(file_path: Path) -> Optional[str]:
    """
    Extract metadata from image files.

    Uses Pillow to get dimensions, format, and basic EXIF data.

    Args:
        file_path: Path to the image file

    Returns:
        Metadata string or None on error
    """
    _init_lazy_imports()
    
    if Image is None or not file_path.exists():
        return None
    
    try:
        with Image.open(file_path) as img:
            width, height = img.size
            format_name = img.format or "Unknown"
            mode = img.mode
            
            metadata_parts = [
                f"Image: {format_name}",
                f"Dimensions: {width}x{height} pixels",
                f"Mode: {mode}",
            ]
            
            # Try to get EXIF data
            try:
                exif = img._getexif()
                if exif:
                    # Get common EXIF tags
                    from PIL.ExifTags import TAGS
                    for tag_id, value in list(exif.items())[:10]:
                        tag = TAGS.get(tag_id, tag_id)
                        if isinstance(value, (str, int, float)):
                            metadata_parts.append(f"{tag}: {value}")
            except (AttributeError, KeyError):
                pass
            
            return "\n".join(metadata_parts)
    except Exception as e:
        logger.warning(f"Failed to extract image metadata from {file_path}: {e}")
        return None


def extract_audio_metadata(file_path: Path) -> Optional[str]:
    """
    Extract metadata from audio files.

    Uses mutagen library to get duration, bitrate, channels, and tags.

    Args:
        file_path: Path to the audio file

    Returns:
        Metadata string or None on error
    """
    _init_lazy_imports()
    
    if mutagen is None or not file_path.exists():
        return None
    
    try:
        audio = mutagen.File(file_path)
        if audio is None:
            return None
        
        metadata_parts = []
        
        # Get basic info
        info = audio.info if hasattr(audio, 'info') else None
        
        if info:
            # Duration
            if hasattr(info, 'length') and info.length:
                duration = info.length
                minutes = int(duration // 60)
                seconds = int(duration % 60)
                metadata_parts.append(f"Duration: {minutes}m {seconds}s")
            
            # Bitrate
            if hasattr(info, 'bitrate') and info.bitrate:
                metadata_parts.append(f"Bitrate: {info.bitrate // 1000} kbps")
            
            # Sample rate
            if hasattr(info, 'sample_rate') and info.sample_rate:
                metadata_parts.append(f"Sample Rate: {info.sample_rate} Hz")
            
            # Channels
            if hasattr(info, 'channels') and info.channels:
                channels = "Stereo" if info.channels == 2 else f"{info.channels} channels"
                metadata_parts.append(f"Channels: {channels}")
        
        # Get format
        format_name = type(audio).__name__.replace(".", " ")
        metadata_parts.insert(0, f"Audio Format: {format_name}")
        
        # Get tags (artist, title, album, etc.)
        tags_to_check = [
            ('TIT2', 'Title'),      # ID3 title
            ('TPE1', 'Artist'),     # ID3 artist
            ('TALB', 'Album'),      # ID3 album
            ('TDRC', 'Year'),       # ID3 year
            ('TCON', 'Genre'),      # ID3 genre
            ('title', 'Title'),     # Vorbis/FLAC
            ('artist', 'Artist'),
            ('album', 'Album'),
            ('date', 'Year'),
            ('genre', 'Genre'),
        ]
        
        for tag_key, tag_name in tags_to_check:
            if tag_key in audio:
                value = audio[tag_key]
                if isinstance(value, list) and value:
                    value = value[0]
                if value:
                    # Truncate long values
                    str_value = str(value)[:100]
                    metadata_parts.append(f"{tag_name}: {str_value}")
        
        return "\n".join(metadata_parts) if metadata_parts else None
    except Exception as e:
        logger.warning(f"Failed to extract audio metadata from {file_path}: {e}")
        return None


def extract_video_metadata(file_path: Path) -> Optional[str]:
    """
    Extract metadata from video files.

    Uses ffprobe (from FFmpeg) to get resolution, duration, codec, etc.

    Args:
        file_path: Path to the video file

    Returns:
        Metadata string or None on error
    """
    _init_lazy_imports()
    
    if ffprobe_path is None or not file_path.exists():
        return None
    
    import subprocess
    import json as json_module
    
    try:
        # Run ffprobe to get video metadata
        cmd = [
            ffprobe_path,
            "-v", "quiet",
            "-print_format", "json",
            "-show_format",
            "-show_streams",
            str(file_path)
        ]
        
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=30
        )
        
        if result.returncode != 0:
            logger.warning(f"ffprobe failed for {file_path}: {result.stderr}")
            return None
        
        data = json_module.loads(result.stdout)
        metadata_parts = []
        
        # Get format info
        format_info = data.get("format", {})
        
        # Duration
        if "duration" in format_info:
            duration = float(format_info["duration"])
            hours = int(duration // 3600)
            minutes = int((duration % 3600) // 60)
            seconds = int(duration % 60)
            if hours > 0:
                metadata_parts.append(f"Duration: {hours}h {minutes}m {seconds}s")
            else:
                metadata_parts.append(f"Duration: {minutes}m {seconds}s")
        
        # File size (already have, but can show bitrate)
        if "bit_rate" in format_info:
            bitrate = int(format_info["bit_rate"]) // 1000
            metadata_parts.append(f"Bitrate: {bitrate} kbps")
        
        # Format name
        if "format_long_name" in format_info:
            metadata_parts.insert(0, f"Video Format: {format_info['format_long_name']}")
        
        # Process streams
        video_stream = None
        audio_stream = None
        
        for stream in data.get("streams", []):
            if stream.get("codec_type") == "video" and not video_stream:
                video_stream = stream
            elif stream.get("codec_type") == "audio" and not audio_stream:
                audio_stream = stream
        
        # Video stream info
        if video_stream:
            width = video_stream.get("width", 0)
            height = video_stream.get("height", 0)
            if width and height:
                # Determine resolution label
                resolution_label = ""
                if height >= 2160:
                    resolution_label = " (4K)"
                elif height >= 1080:
                    resolution_label = " (1080p)"
                elif height >= 720:
                    resolution_label = " (720p)"
                elif height >= 480:
                    resolution_label = " (480p)"
                metadata_parts.append(f"Resolution: {width}x{height}{resolution_label}")
            
            # Video codec
            codec = video_stream.get("codec_long_name") or video_stream.get("codec_name")
            if codec:
                metadata_parts.append(f"Video Codec: {codec}")
            
            # Frame rate
            fps_str = video_stream.get("r_frame_rate", "")
            if fps_str and "/" in fps_str:
                num, den = fps_str.split("/")
                if int(den) > 0:
                    fps = int(num) / int(den)
                    metadata_parts.append(f"Frame Rate: {fps:.2f} fps")
        
        # Audio stream info
        if audio_stream:
            audio_codec = audio_stream.get("codec_long_name") or audio_stream.get("codec_name")
            if audio_codec:
                metadata_parts.append(f"Audio Codec: {audio_codec}")
            
            channels = audio_stream.get("channels")
            if channels:
                channel_str = "Stereo" if channels == 2 else f"{channels} channels"
                metadata_parts.append(f"Audio: {channel_str}")
        
        # Tags (title, etc.)
        tags = format_info.get("tags", {})
        for tag_key in ["title", "artist", "album", "date", "comment"]:
            if tag_key in tags:
                value = str(tags[tag_key])[:100]
                metadata_parts.append(f"{tag_key.capitalize()}: {value}")
        
        return "\n".join(metadata_parts) if metadata_parts else None
    except subprocess.TimeoutExpired:
        logger.warning(f"ffprobe timeout for {file_path}")
        return None
    except Exception as e:
        logger.warning(f"Failed to extract video metadata from {file_path}: {e}")
        return None


# =============================================================================
# Extractor Class
# =============================================================================

class Extractor:
    """
    Content extractor for various file formats.

    Enriches FileRecord objects with content_excerpt and mime fields.

    Attributes:
        max_excerpt_bytes: Maximum excerpt size in bytes
        max_pdf_pages: Maximum PDF pages to extract
        stats: Dictionary tracking extraction statistics
    """

    def __init__(
        self,
        max_excerpt_bytes: int = DEFAULT_MAX_EXCERPT_BYTES,
        max_pdf_pages: int = DEFAULT_MAX_PDF_PAGES,
    ):
        """
        Initialize Extractor with configuration.

        Args:
            max_excerpt_bytes: Maximum excerpt size (default 8KB)
            max_pdf_pages: Maximum PDF pages to extract (default 5)
        """
        self.max_excerpt_bytes = max_excerpt_bytes
        self.max_pdf_pages = max_pdf_pages
        
        # Statistics tracking
        self.stats = {
            "files_processed": 0,
            "extraction_errors": 0,
            "total_excerpt_bytes": 0,
        }

    def _reset_stats(self) -> None:
        """Reset statistics counters."""
        self.stats = {
            "files_processed": 0,
            "extraction_errors": 0,
            "total_excerpt_bytes": 0,
        }

    def _extract_content(self, file_path: Path, extension: str) -> Optional[str]:
        """
        Extract content based on file type.

        Args:
            file_path: Path to the file
            extension: File extension (lowercase)

        Returns:
            Extracted content or None
        """
        ext = extension.lower()
        
        # Plain text files
        if ext in TEXT_EXTENSIONS:
            return extract_text_content(file_path, self.max_excerpt_bytes)
        
        # PDF files
        if ext == ".pdf":
            return extract_pdf_content(
                file_path,
                max_pages=self.max_pdf_pages,
                max_bytes=self.max_excerpt_bytes
            )
        
        # Word documents
        if ext == ".docx":
            return extract_docx_content(file_path, self.max_excerpt_bytes)
        
        if ext == ".doc":
            # Legacy .doc format not supported
            return None
        
        # PowerPoint
        if ext == ".pptx":
            return extract_pptx_content(file_path, self.max_excerpt_bytes)
        
        # Excel
        if ext in (".xlsx", ".xls"):
            return extract_xlsx_content(file_path, max_bytes=self.max_excerpt_bytes)
        
        # Images (metadata only)
        if ext in IMAGE_EXTENSIONS:
            return extract_image_metadata(file_path)
        
        # Audio files (metadata only)
        if ext in AUDIO_EXTENSIONS:
            return extract_audio_metadata(file_path)
        
        # Video files (metadata only)
        if ext in VIDEO_EXTENSIONS:
            return extract_video_metadata(file_path)
        
        # Unsupported format
        return None

    def extract(self, record: FileRecord) -> FileRecord:
        """
        Extract content and enrich FileRecord.

        Args:
            record: FileRecord to enrich

        Returns:
            New FileRecord with content_excerpt and mime fields
        """
        self.stats["files_processed"] += 1
        
        # Detect MIME type
        mime = detect_mime_type(record.path)
        
        # Extract content
        try:
            content = self._extract_content(record.path, record.extension)
            
            if content:
                self.stats["total_excerpt_bytes"] += len(content.encode("utf-8"))
        except Exception as e:
            logger.warning(f"Extraction error for {record.path}: {e}")
            self.stats["extraction_errors"] += 1
            content = None
        
        # Create enriched record (immutable pattern)
        return FileRecord(
            path=record.path,
            size=record.size,
            mtime=record.mtime,
            ctime=record.ctime,
            sha256=record.sha256,
            extension=record.extension,
            mime=mime,
            content_excerpt=content,
        )

    def extract_batch(
        self,
        records: List[FileRecord],
        callback=None
    ) -> Generator[FileRecord, None, None]:
        """
        Extract content for multiple FileRecords.

        Args:
            records: List of FileRecords to process
            callback: Optional callback(count, path) for progress

        Yields:
            Enriched FileRecord objects
        """
        for i, record in enumerate(records):
            enriched = self.extract(record)
            
            if callback:
                callback(i + 1, record.path)
            
            yield enriched
